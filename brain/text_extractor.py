"""
Text Extractor — Extract plain text from PDF, DOCX, PPTX, MD, and TXT files.

Returns: { content: str, error: str | None, metadata: dict }
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Optional


def extract_text(file_path: str) -> dict:
    """
    Extract plain text from a file based on its extension.

    Returns dict with keys: content, error, metadata.
    """
    path = Path(file_path)
    if not path.exists():
        return {"content": "", "error": f"File not found: {file_path}", "metadata": {}}

    ext = path.suffix.lower()
    size = path.stat().st_size

    metadata = {
        "file_name": path.name,
        "file_size": size,
        "file_type": ext.lstrip("."),
    }

    try:
        if ext == ".pdf":
            content = _extract_pdf(path)
        elif ext == ".docx":
            content = _extract_docx(path)
        elif ext == ".pptx":
            content = _extract_pptx(path)
        elif ext in (".md", ".txt", ".text", ".markdown"):
            content = path.read_text(encoding="utf-8", errors="replace")
        else:
            return {
                "content": "",
                "error": f"Unsupported file type: {ext}",
                "metadata": metadata,
            }

        metadata["char_count"] = len(content)
        return {"content": content, "error": None, "metadata": metadata}

    except Exception as e:
        return {"content": "", "error": str(e), "metadata": metadata}


def _extract_pdf(path: Path) -> str:
    """Extract text from PDF using pdfplumber."""
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            # Also extract table text
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        cells = [str(c) if c else "" for c in row]
                        text += "\n" + " | ".join(cells)
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    import docx

    doc = docx.Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides):
        slide_text = [f"--- Slide {i + 1} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(f"[Notes: {notes}]")
        parts.append("\n".join(slide_text))

    return "\n\n".join(parts)


def get_file_type(filename: str) -> Optional[str]:
    """Return normalized file type from filename, or None if unsupported."""
    ext = Path(filename).suffix.lower().lstrip(".")
    mapping = {
        "pdf": "pdf",
        "docx": "docx",
        "pptx": "pptx",
        "md": "md",
        "markdown": "md",
        "txt": "txt",
        "text": "txt",
    }
    return mapping.get(ext)


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".txt", ".text", ".markdown"}
